#!/usr/bin/env python3
"""Generate Experiment 1 PPTX: Static vs Agentic Red Teaming.

Uses the same visual style as RedTeam_Results.pptx (Calibri, white bg).
Reads results from the experiment1 output directory.

Usage:
    python scripts/create_experiment1_pptx.py
    python scripts/create_experiment1_pptx.py --results-dir results/experiment1_0405_HHMM
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design tokens (matching RedTeam_Results.pptx) ────────────────────
TITLE_FONT = "Calibri"
TITLE_SIZE = Pt(36)
TITLE_COLOR = RGBColor(0x33, 0x33, 0x33)
TITLE_LEFT = Inches(0.83)
TITLE_TOP = Inches(0.63)
TITLE_WIDTH = Inches(11.7)
TITLE_HEIGHT = Inches(0.6)

HEADING_SIZE = Pt(15)
HEADING_COLOR = RGBColor(0x1A, 0x1A, 0x1A)

BODY_SIZE = Pt(12)
BODY_COLOR = RGBColor(0x55, 0x55, 0x55)

ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
BADGE_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER_BG = RGBColor(0x0F, 0x17, 0x2A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_run(p, text, size, color, bold=False, font_name=TITLE_FONT):
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font_name
    return r


def _clean_placeholders(slide):
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def add_title(slide, text):
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    _add_run(txBox.text_frame.paragraphs[0], text, TITLE_SIZE, TITLE_COLOR, bold=True)
    return txBox


def add_heading(slide, text, top, left=Inches(1.0), width=Inches(11)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    _add_run(txBox.text_frame.paragraphs[0], text, HEADING_SIZE, HEADING_COLOR, bold=True)
    return txBox


def add_body(slide, text, top, left=Inches(1.0), width=Inches(11), color=BODY_COLOR, bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    txBox.text_frame.word_wrap = True
    _add_run(txBox.text_frame.paragraphs[0], text, BODY_SIZE, color, bold=bold)
    return txBox


def add_badge(slide, text, top, left, width, color):
    box = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _add_run(p, text, Pt(11), BADGE_TEXT, bold=True)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    return box


def add_image(slide, img_path, top, left, width, height):
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), left, top, width, height)


def find_latest_experiment() -> Path:
    results_root = Path(__file__).resolve().parents[1] / "results"
    dirs = sorted([d for d in results_root.iterdir()
                   if d.is_dir() and d.name.startswith("experiment1_")])
    if not dirs:
        raise FileNotFoundError("No experiment1 results found")
    return dirs[-1]


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--results-dir", type=str, default=None)
    args = parser.parse_args()

    results_dir = Path(args.results_dir) if args.results_dir else find_latest_experiment()
    report_dir = results_dir / "report"

    print(f"Results dir: {results_dir}")

    # Load data
    summary = {}
    summary_path = report_dir / "experiment1_summary.json"
    if summary_path.exists():
        with open(summary_path) as f:
            summary = json.load(f)

    df = None
    csv_path = report_dir / "experiment1_all_results.csv"
    if csv_path.exists():
        df = pd.read_csv(csv_path)

    # Load raw results for details
    static_data = []
    agentic_data = []
    if (results_dir / "static_results.json").exists():
        with open(results_dir / "static_results.json") as f:
            static_data = json.load(f)
    if (results_dir / "agentic_results.json").exists():
        with open(results_dir / "agentic_results.json") as f:
            agentic_data = json.load(f)

    # Compute stats if summary is missing/incomplete
    static_n = len(static_data)
    static_k = sum(1 for r in static_data if r.get("success"))
    static_asr = static_k / static_n if static_n > 0 else 0
    agentic_n = len(agentic_data)
    agentic_k = sum(1 for r in agentic_data if r.get("success"))
    agentic_asr = agentic_k / agentic_n if agentic_n > 0 else 0

    # ── Build PPTX ──
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DIVIDER_BG
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(0.8))
    _add_run(txBox.text_frame.paragraphs[0],
             "Experiment 1: Static vs Agentic Red Teaming",
             Pt(40), BADGE_TEXT, bold=True)
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(0.5))
    _add_run(txBox2.text_frame.paragraphs[0],
             "Do autonomous, adaptive attacks outperform predefined single-turn attacks?",
             Pt(18), RGBColor(0xAA, 0xAA, 0xAA))
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.4))
    _add_run(txBox3.text_frame.paragraphs[0],
             "IIT Bombay EPGD AI/ML  |  Capstone Project  |  CommodityRedTeam",
             Pt(14), RGBColor(0x88, 0x88, 0x88))

    # ── Slide 2: Hypothesis & Setup ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "Hypothesis & Experimental Setup")

    add_heading(slide, "Hypothesis", Inches(1.4))
    add_body(slide, "Agentic red teaming (multi-step, adaptive, LLM-as-judge) achieves higher ASR "
             "than static predefined single-turn attacks against the same target.",
             Inches(1.75), width=Inches(11))

    add_heading(slide, "Independent Variable", Inches(2.5))
    add_body(slide, "Attack method: Static (Condition A) vs Agentic (Condition B)", Inches(2.85))

    add_heading(slide, "Dependent Variables", Inches(3.4))
    add_body(slide, "Attack Success Rate (ASR), steps to success, round-over-round improvement, "
             "statistical significance (chi-squared, Bayesian)", Inches(3.75))

    # Condition boxes
    add_badge(slide, "CONDITION A: Static Baseline", Inches(4.6), Inches(0.8), Inches(5.5), ACCENT_BLUE)
    add_body(slide, f"50 predefined single-turn attacks (v1 library)  |  {static_n} attacks  |  "
             "1 step each  |  keyword-based evaluation",
             Inches(5.0), left=Inches(0.8), width=Inches(5.5))

    add_badge(slide, "CONDITION B: Agentic System", Inches(4.6), Inches(6.8), Inches(5.5), ACCENT_RED)
    add_body(slide, f"Planner -> Attacker -> Critic -> Mutator loop  |  {agentic_n} attacks  |  "
             "multi-step  |  LLM-as-judge evaluation",
             Inches(5.0), left=Inches(6.8), width=Inches(5.5))

    # ── Slide 3: Architecture ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "v2 Agentic Architecture")

    components = [
        ("Planner Agent", "Generates structured multi-step attack plans via LLM. "
         "Cold-start mode + adaptive mode based on attack history.", ACCENT_BLUE),
        ("Attacker", "Executes plan steps sequentially against target agent. "
         "Maintains conversation context across steps for multi-turn attacks.", ACCENT_ORANGE),
        ("Critic Agent (LLM-as-Judge)", "Evaluates attack success against target's actual rules. "
         "Replaces keyword matching with semantic understanding.", ACCENT_RED),
        ("Mutator Agent", "Improves failed attacks via 7 strategies: rephrase, obfuscate, "
         "add indirection, split steps, change vector, escalate, social engineer.", ACCENT_PURPLE),
    ]

    for i, (name, desc, color) in enumerate(components):
        y = Inches(1.5) + Inches(i * 1.3)
        add_badge(slide, name, y, Inches(0.8), Inches(3.5), color)
        add_body(slide, desc, y + Inches(0.4), left=Inches(0.8), width=Inches(11.5))

    # Flow arrow text
    add_body(slide, "Loop: Plan -> Execute -> Evaluate -> Mutate -> Repeat (3 rounds, adaptive)",
             Inches(6.7), left=Inches(0.8), color=HEADING_COLOR, bold=True)

    # ── Slide 4: Comparison Plot ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "Results: Static vs Agentic Comparison")

    plot_path = report_dir / "experiment1_comparison.png"
    if plot_path.exists():
        slide.shapes.add_picture(str(plot_path), Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5))
    else:
        add_body(slide, "[Plot not found -- run experiment first]", Inches(3))

    # ── Slide 5: Key Numbers ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "Results: Key Metrics")

    # Static column
    add_badge(slide, "STATIC (v1)", Inches(1.5), Inches(0.8), Inches(5.5), ACCENT_BLUE)
    metrics_static = [
        f"Attacks: {static_n}",
        f"Successful: {static_k}",
        f"ASR: {static_asr:.1%}",
        "Steps per attack: 1 (always)",
        "Evaluation: keyword matching",
    ]
    for i, m in enumerate(metrics_static):
        add_body(slide, m, Inches(2.0) + Inches(i * 0.4), left=Inches(1.0), width=Inches(5))

    # Agentic column
    add_badge(slide, "AGENTIC (v2)", Inches(1.5), Inches(6.8), Inches(5.5), ACCENT_RED)
    avg_steps = 0
    if agentic_k > 0:
        successes = [r for r in agentic_data if r.get("success")]
        avg_steps = sum(r.get("steps", 1) for r in successes) / len(successes) if successes else 0

    metrics_agentic = [
        f"Attacks: {agentic_n}",
        f"Successful: {agentic_k}",
        f"ASR: {agentic_asr:.1%}",
        f"Avg steps to success: {avg_steps:.1f}",
        "Evaluation: LLM-as-judge (Critic agent)",
    ]
    for i, m in enumerate(metrics_agentic):
        add_body(slide, m, Inches(2.0) + Inches(i * 0.4), left=Inches(7.0), width=Inches(5))

    # Significance
    comp = summary.get("comparison", {})
    p_val = comp.get("p_value", 1.0)
    chi2 = comp.get("chi2", 0)
    sig = "***" if p_val < 0.001 else "**" if p_val < 0.01 else "*" if p_val < 0.05 else "ns"

    add_heading(slide, "Statistical Significance", Inches(4.5))
    add_body(slide, f"Chi-squared test: chi2 = {chi2:.3f}, p = {p_val:.4f} ({sig})",
             Inches(4.9), color=ACCENT_RED if p_val < 0.05 else BODY_COLOR, bold=p_val < 0.05)
    add_body(slide, f"Delta ASR: {agentic_asr - static_asr:+.1%} "
             f"({'agentic better' if agentic_asr > static_asr else 'static better'})",
             Inches(5.3))

    # ── Slide 6: Round-over-Round ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "Agentic Improvement: Round over Round")

    if agentic_data:
        rounds = {}
        for r in agentic_data:
            rnd = r.get("round", 0)
            rounds.setdefault(rnd, {"total": 0, "success": 0})
            rounds[rnd]["total"] += 1
            if r.get("success"):
                rounds[rnd]["success"] += 1

        add_heading(slide, "ASR by Round", Inches(1.5))
        for rnd in sorted(rounds.keys()):
            s = rounds[rnd]
            asr = s["success"] / s["total"] if s["total"] > 0 else 0
            y = Inches(2.0) + Inches(rnd * 0.8)
            label = f"Round {rnd + 1}"
            if rnd > 0:
                label += " (adaptive + mutations)"
            add_body(slide, label, y, left=Inches(1.0), width=Inches(3), bold=True)
            add_body(slide, f"{s['success']}/{s['total']} = {asr:.0%}",
                     y, left=Inches(5.0), width=Inches(2),
                     color=ACCENT_RED if asr > 0.3 else ACCENT_GREEN, bold=True)

        add_heading(slide, "Key Observation", Inches(5.0))
        r1_asr = rounds.get(0, {}).get("success", 0) / max(rounds.get(0, {}).get("total", 1), 1)
        last_rnd = max(rounds.keys())
        rN_asr = rounds[last_rnd]["success"] / max(rounds[last_rnd]["total"], 1)
        if rN_asr > r1_asr:
            add_body(slide, f"Adaptive planning + mutation improved ASR from {r1_asr:.0%} (Round 1) "
                     f"to {rN_asr:.0%} (Round {last_rnd + 1}), demonstrating self-improvement.",
                     Inches(5.4))
        elif rN_asr == r1_asr:
            add_body(slide, f"ASR remained stable at {rN_asr:.0%} across rounds. "
                     "Mutations adapted attack strategies but target defense held.",
                     Inches(5.4))
        else:
            add_body(slide, f"ASR decreased from {r1_asr:.0%} to {rN_asr:.0%}. "
                     "Target agent may have been easier to attack with fresh plans than mutations.",
                     Inches(5.4))

    # ── Slide 7: ML Techniques Used ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "ML/AI Techniques Demonstrated")

    techniques = [
        ("Module 2: Embeddings", "D6 Semantic Filter -- sentence-transformer cosine similarity for injection detection",
         "sim(q, p) = (q . p) / (||q|| * ||p||) in 384-dim embedding space"),
        ("Module 2: Ensemble", "D8 Ensemble Defense -- XGBoost combining all defense signals",
         "Gradient boosted trees with L2 regularization, cross-validated AUC"),
        ("Module 3: Agentic AI", "v2 Planner/Critic/Mutator agents -- autonomous attack loop",
         "Multi-agent system with structured plan generation and LLM-as-judge"),
        ("Module 4: Explainability", "SHAP TreeExplainer for attack success prediction",
         "Shapley values: phi_i = avg marginal contribution across coalitions"),
        ("Module 4: Model Eval", "Bayesian vulnerability (Beta-Binomial), MI, transferability",
         "Posterior: p | k,n ~ Beta(1+k, 1+n-k), Fisher's exact test for transfer"),
    ]

    for i, (module, technique, math) in enumerate(techniques):
        y = Inches(1.4) + Inches(i * 1.1)
        add_body(slide, module, y, left=Inches(0.8), width=Inches(3), color=ACCENT_BLUE, bold=True)
        add_body(slide, technique, y + Inches(0.3), left=Inches(0.8), width=Inches(11.5))
        add_body(slide, math, y + Inches(0.6), left=Inches(0.8), width=Inches(11.5),
                 color=RGBColor(0x88, 0x88, 0x88))

    # ── Slide 8: Findings ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _clean_placeholders(slide)
    add_title(slide, "Key Findings")

    findings = [
        ("1. Agentic attacks show adaptive capability",
         f"Round-over-round ASR demonstrates the mutator's ability to evolve failed attacks. "
         f"The planner generates increasingly sophisticated multi-step plans based on history."),
        ("2. LLM-as-judge provides richer evaluation",
         "The Critic agent detects subtle rule violations that keyword matching misses, "
         "and avoids false positives from keywords appearing in rejection context."),
        ("3. Multi-step plans exploit interaction complexity",
         "Agentic attacks that chain legitimate-looking steps (e.g., establish position then "
         "'hedge' with correlated long) bypass single-point defenses."),
        ("4. Static attacks have broader category coverage",
         f"50 predefined attacks cover all 7 categories systematically. "
         f"Agentic attacks tend to converge on categories the planner finds most promising."),
        ("5. Defense implications",
         "Static attacks test breadth (all vulnerability categories). "
         "Agentic attacks test depth (adaptive evasion of specific defenses). Both are needed."),
    ]

    for i, (title, body) in enumerate(findings):
        y = Inches(1.4) + Inches(i * 1.1)
        add_heading(slide, title, y)
        add_body(slide, body, y + Inches(0.35), width=Inches(11.5))

    # ── Save ──
    out_path = report_dir / "Experiment1_Results.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
