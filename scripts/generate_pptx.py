#!/usr/bin/env python3
"""Generate PPTX report from benchmark results using the IIT Bombay template.

Usage:
    python scripts/generate_pptx.py
    python scripts/generate_pptx.py --results-dir results/results_0329_1945
    python scripts/generate_pptx.py --results-dir results/results_0329_1945 --template ~/Downloads/AI_ML_Sample_Format.pptx
"""

from __future__ import annotations

import argparse
import csv
import sys
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt


# ── Helpers ──────────────────────────────────────────────────────


def add_textbox(
    slide, left, top, width, height, text,
    font_size=14, bold=False, color=RGBColor(0x33, 0x33, 0x33),
    alignment=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_title(slide, text):
    add_textbox(
        slide, 762000, 571500, 11201400, 554100, text,
        font_size=36, bold=True, color=RGBColor(0x33, 0x33, 0x33),
    )


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])


def add_image(slide, img_path, left, top, width, height):
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Emu(left), Emu(top), Emu(width), Emu(height))


def add_table(slide, left, top, width, rows_data, col_widths=None):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_height = Emu(280000 * n_rows)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, Emu(left), Emu(top), Emu(width), tbl_height,
    )
    table = table_shape.table
    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(10)
                para.font.name = "Calibri"
                if r_idx == 0:
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    para.font.size = Pt(11)
                else:
                    para.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            if r_idx == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", "2D2D2D")
            elif r_idx % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
                srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
                srgbClr.set("val", "F5F5F5")
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Emu(w)
    return table_shape


def load_summary(results_dir: Path) -> list[dict]:
    """Load summary.csv into a list of dicts."""
    summary_path = results_dir / "summary.csv"
    if not summary_path.exists():
        return []
    with open(summary_path) as f:
        return list(csv.DictReader(f))


def build_table_data(summary: list[dict]) -> list[list[str]]:
    """Convert summary rows into table data for the PPTX."""
    header = ["Model", "Defense", "ASR %", "Detection %", "Impact ($M)"]
    rows = [header]
    for r in summary:
        impact_m = round(int(r["total_impact_usd"]) / 1_000_000, 1)
        rows.append([
            r["model"], r["defense"], r["asr_pct"],
            r["detection_rate_pct"], str(impact_m),
        ])
    return rows


def build_ranking(summary: list[dict], model: str) -> list[list[str]]:
    """Build defense ranking table for a given model."""
    model_rows = [r for r in summary if r["model"] == model]
    baseline = next((r for r in model_rows if r["defense"] == "none"), None)
    baseline_asr = float(baseline["asr_pct"]) if baseline else 0

    ranked = []
    for r in model_rows:
        if r["defense"] == "none" or "+" in r["defense"]:
            continue
        asr = float(r["asr_pct"])
        diff = asr - baseline_asr
        sign = "+" if diff >= 0 else ""
        ranked.append((r["defense"], asr, f"{sign}{diff:.0f}pp"))

    ranked.sort(key=lambda x: x[1])

    header = ["Rank", "Defense", "ASR", "vs Baseline"]
    rows = [header]
    for i, (defense, asr, diff) in enumerate(ranked, 1):
        rows.append([str(i), defense, f"{asr:.0f}%", diff])
    return rows


# ── Slide builders ───────────────────────────────────────────────


def fill_template_slides(prs):
    """Fill slides 1-7 (the template's project context slides)."""
    slides = list(prs.slides)

    # Slide 1: Title
    for shape in slides[0].shapes:
        if shape.has_text_frame:
            t = shape.text_frame.paragraphs[0].text
            if "AI/ML in Practice" in t:
                shape.text_frame.paragraphs[0].text = "AI/ML in Practice"
            elif "Title" == t.strip():
                shape.text_frame.paragraphs[0].text = "Red Teaming Agentic AI"
            elif "Your Names" in t or "Sample" in t:
                shape.text_frame.paragraphs[0].text = "Capstone Project: IIT Bombay EPGD AI/ML"

    # Slide 2: Problem Statement and Motivation (The "Why")
    card_titles, card_descs = [], []
    for shape in slides[1].shapes:
        if shape.has_text_frame and shape.top > 2500000:
            runs = shape.text_frame.paragraphs[0].runs
            fsize = runs[0].font.size if runs else shape.text_frame.paragraphs[0].font.size
            if fsize and fsize >= Pt(18):
                card_titles.append(shape)
            elif fsize and fsize < Pt(18):
                card_descs.append(shape)

    card_titles.sort(key=lambda s: s.left)
    card_descs.sort(key=lambda s: s.left)

    for i, t in enumerate(["Agentic AI in Finance", "No Systematic Testing", "Real Financial Risk"]):
        if i < len(card_titles):
            card_titles[i].text_frame.paragraphs[0].text = t
    for i, d in enumerate([
        "LLM-based agents are making autonomous trading decisions on commodity markets worth trillions",
        "No established framework exists to red-team domain-specific agentic AI systems for vulnerabilities",
        "A single manipulated trade recommendation can cause millions in losses -- 5% price error cascades to $10M+ impact",
    ]):
        if i < len(card_descs):
            card_descs[i].text_frame.paragraphs[0].text = d

    # Slide 3: Proposed Solution (The "What")
    bullets3 = sorted(
        [s for s in slides[2].shapes if s.has_text_frame and s.left < 5500000 and s.top > 1000000],
        key=lambda s: s.top,
    )
    for i, t in enumerate([
        "CommodityRedTeam: A research framework with 50 domain-specific attacks across 7 vulnerability categories (prompt injection, tool manipulation, reasoning hijacking, etc.)",
        "5 layered defense strategies: input filtering, output validation, system prompt hardening, multi-agent verification, and human-in-the-loop review",
        "Automated evaluation pipeline with ASR, detection rate, financial impact metrics, and statistical significance testing (chi-squared, Cohen's h)",
    ]):
        if i < len(bullets3):
            bullets3[i].text_frame.paragraphs[0].text = t
            bullets3[i].text_frame.paragraphs[0].font.size = Pt(13)

    # Slide 4: Data Sources and Tools (The "How")
    bullets4 = sorted(
        [s for s in slides[3].shapes if s.has_text_frame and s.top > 1000000],
        key=lambda s: s.top,
    )
    for i, t in enumerate([
        "LLMs: Groq LPU (Llama 3.3 70B, Qwen3 32B, Llama 4 Scout), Google Gemini, Mistral -- all free-tier APIs via OpenAI-compatible endpoints",
        "Tools & Frameworks: Python, LangChain (agent + tool calling), yfinance (real market data), Rich (CLI), Matplotlib/Seaborn (viz), scipy/statsmodels (statistics)",
        "Target Agent: Commodity trading agent with 7 tools (price, news, risk, fundamentals, correlation, position limits, recommendation) -- each with switchable attack modes",
    ]):
        if i < len(bullets4):
            bullets4[i].text_frame.paragraphs[0].text = t
            bullets4[i].text_frame.paragraphs[0].font.size = Pt(13)

    # Slide 5: Roadmap (The "When")
    step_titles = sorted(
        [s for s in slides[4].shapes if s.has_text_frame and "Step" in s.text_frame.paragraphs[0].text],
        key=lambda s: s.left,
    )
    step_descs = sorted(
        [s for s in slides[4].shapes if s.has_text_frame and s.text_frame.paragraphs[0].text in ("\u2014", "\u2014")],
        key=lambda s: s.left,
    )
    for i, t in enumerate(["Build Agent", "Design Attacks", "Implement Defenses", "Benchmark & Report"]):
        if i < len(step_titles):
            step_titles[i].text_frame.paragraphs[0].text = t
    for i, d in enumerate([
        "LLM trading agent\nwith 7 tools", "50 attacks across\n7 categories",
        "5 defense strategies\nlayered approach", "Automated pipeline\nstatistical analysis",
    ]):
        if i < len(step_descs):
            step_descs[i].text_frame.paragraphs[0].text = d

    # Slide 6: Key Deliverables
    bullets6 = sorted(
        [s for s in slides[5].shapes if s.has_text_frame and s.left < 5500000 and s.top > 1000000],
        key=lambda s: s.top,
    )
    for i, t in enumerate([
        "CommodityAgentThreat Taxonomy: First vulnerability classification for LLM-based financial trading agents (7 categories, 50 attacks)",
        "Cross-model vulnerability comparison: Qwen3 32B vs Llama 4 Scout -- different models need different defense strategies",
        "Empirical defense evaluation: Layered defenses reduce ASR from 38-42% to 0-12% with statistical significance",
    ]):
        if i < len(bullets6):
            bullets6[i].text_frame.paragraphs[0].text = t
            bullets6[i].text_frame.paragraphs[0].font.size = Pt(13)

    # Slide 7: Discussion -- keep as-is


def add_results_slides(prs, results_dir: Path):
    """Add slides 8-14 with benchmark results."""
    report_dir = results_dir / "report"
    summary = load_summary(results_dir)

    if not summary:
        print("Warning: No summary.csv found, skipping results slides.")
        return

    models = sorted(set(r["model"] for r in summary))

    # Slide 8: Results table
    s8 = add_blank_slide(prs)
    add_title(s8, "Benchmark Results: ASR by Defense Configuration")
    table_data = build_table_data(summary)
    add_table(
        s8, 762000, 1400000, 10668000, table_data,
        col_widths=[2200000, 2400000, 1600000, 2000000, 2400000],
    )

    # Slide 9: ASR Heatmap
    s9 = add_blank_slide(prs)
    add_title(s9, "Attack Success Rate by Category")
    add_image(s9, report_dir / "heatmap_asr.png", 1200000, 1400000, 9800000, 5000000)

    # Slide 10: Defense Bar Chart
    s10 = add_blank_slide(prs)
    add_title(s10, "Defense Effectiveness Comparison")
    add_image(s10, report_dir / "barchart_defense_asr.png", 762000, 1400000, 10668000, 5000000)

    # Slide 11: Radar Chart
    s11 = add_blank_slide(prs)
    add_title(s11, "Model Vulnerability Profiles")
    add_image(s11, report_dir / "radar_vulnerability.png", 1500000, 1300000, 9200000, 5200000)

    # Slide 12: Detection Coverage
    s12 = add_blank_slide(prs)
    add_title(s12, "Detection Coverage by Defense")
    add_image(s12, report_dir / "heatmap_detection_coverage.png", 1200000, 1400000, 9800000, 5000000)

    # Slide 13: Defense Rankings
    s13 = add_blank_slide(prs)
    add_title(s13, "Defense Effectiveness Ranking")

    x_positions = [762000, 6400000]
    for i, model in enumerate(models[:2]):
        x = x_positions[i]
        add_textbox(s13, x, 1300000, 5000000, 300000, model, font_size=16, bold=True)
        ranking = build_ranking(summary, model)
        add_table(
            s13, x, 1700000, 5000000, ranking,
            col_widths=[800000, 1800000, 1000000, 1400000],
        )

    # Slide 14: Key Findings
    s14 = add_blank_slide(prs)
    add_title(s14, "Key Findings")

    # Build findings dynamically from data
    findings = _derive_findings(summary, models)
    y = 1400000
    for title, desc in findings:
        add_textbox(
            s14, 900000, y, 10400000, 250000, title,
            font_size=15, bold=True, color=RGBColor(0x1A, 0x1A, 0x1A),
        )
        add_textbox(
            s14, 900000, y + 280000, 10400000, 300000, desc,
            font_size=12, color=RGBColor(0x55, 0x55, 0x55),
        )
        y += 700000


def _derive_findings(summary: list[dict], models: list[str]) -> list[tuple[str, str]]:
    """Generate key findings from the summary data."""
    findings = []

    # Find best/worst model with all_combined
    combined = [r for r in summary if "+" in r.get("defense", "")]
    if combined:
        best = min(combined, key=lambda r: float(r["asr_pct"]))
        worst = max(combined, key=lambda r: float(r["asr_pct"]))
        findings.append((
            f"1. {best['model']} is the most robust model",
            f"All defenses combined achieved {best['asr_pct']}% ASR "
            f"(vs {worst['asr_pct']}% for {worst['model']}).",
        ))

    # Best single defense per model
    for model in models:
        model_rows = [r for r in summary if r["model"] == model]
        baseline = next((r for r in model_rows if r["defense"] == "none"), None)
        singles = [r for r in model_rows if r["defense"] != "none" and "+" not in r["defense"]]
        if singles and baseline:
            best_def = min(singles, key=lambda r: float(r["asr_pct"]))
            findings.append((
                f"{len(findings) + 1}. Best defense for {model}: {best_def['defense']}",
                f"ASR reduced from {baseline['asr_pct']}% to {best_def['asr_pct']}%.",
            ))

    # Layered defense impact
    for model in models:
        baseline = next((r for r in summary if r["model"] == model and r["defense"] == "none"), None)
        combined = next((r for r in summary if r["model"] == model and "+" in r.get("defense", "")), None)
        if baseline and combined:
            b_impact = int(baseline["total_impact_usd"])
            c_impact = int(combined["total_impact_usd"])
            if b_impact > 0:
                reduction = round(100 * (1 - c_impact / b_impact))
                findings.append((
                    f"{len(findings) + 1}. Financial impact reduction for {model}: {reduction}%",
                    f"${b_impact / 1e6:.0f}M reduced to ${c_impact / 1e6:.0f}M with all defenses combined.",
                ))

    return findings[:5]  # Cap at 5 findings


# ── Main ─────────────────────────────────────────────────────────


def main():
    parser = argparse.ArgumentParser(description="Generate PPTX report from benchmark results")
    parser.add_argument(
        "--results-dir", type=Path,
        help="Path to results directory (default: latest in results/)",
    )
    parser.add_argument(
        "--template", type=Path,
        default=Path.home() / "winhome" / "Downloads" / "AI_ML_Sample_Format.pptx",
        help="Path to PPTX template",
    )
    parser.add_argument(
        "--output", type=Path, default=None,
        help="Output PPTX path (default: <results-dir>/report/RedTeam_Results.pptx)",
    )
    args = parser.parse_args()

    # Find results dir
    project_dir = Path(__file__).resolve().parents[1]
    if args.results_dir:
        results_dir = args.results_dir
    else:
        results_root = project_dir / "results"
        dirs = sorted(results_root.glob("results_*"), key=lambda p: p.stat().st_mtime, reverse=True)
        if not dirs:
            print("No results directories found.")
            sys.exit(1)
        results_dir = dirs[0]
        print(f"Using latest results: {results_dir}")

    if not args.template.exists():
        print(f"Template not found: {args.template}")
        sys.exit(1)

    output = args.output or (results_dir / "report" / "RedTeam_Results.pptx")
    output.parent.mkdir(parents=True, exist_ok=True)

    # Build presentation
    prs = Presentation(str(args.template))
    fill_template_slides(prs)
    add_results_slides(prs, results_dir)

    prs.save(str(output))
    print(f"Saved: {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
